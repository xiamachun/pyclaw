"""Handle WebSocket connection and messages."""

import json
import uuid
from pathlib import Path
from typing import Any

from fastapi import WebSocket

from pyclaw.constants import (
    DEFAULT_MAX_TOKENS,
    DEFAULT_TEMPERATURE,
    MAX_HISTORY_MESSAGES,
    MAX_SKILL_CONTENT_LENGTH,
)
from pyclaw.sessions.manager import SessionManager
from pyclaw.sessions.transcript import append_message, read_messages


class WebSocketManager:
    """WebSocket connection manager."""

    def __init__(self) -> None:
        """Initialize manager."""
        self.active_connections: list[WebSocket] = []

    async def connect(self, websocket: WebSocket) -> None:
        """Accept new connection.

        Args:
            websocket: WebSocket connection
        """
        await websocket.accept()
        self.active_connections.append(websocket)

    async def disconnect(self, websocket: WebSocket) -> None:
        """Disconnect.

        Args:
            websocket: WebSocket connection
        """
        if websocket in self.active_connections:
            self.active_connections.remove(websocket)

    async def broadcast(self, message: dict[str, Any]) -> None:
        """Broadcast message to all connections.

        Args:
            message: Message content
        """
        for connection in self.active_connections:
            await self.send_to(connection, message)

    async def send_to(self, websocket: WebSocket, message: dict[str, Any]) -> None:
        """Send message to specific connection.

        Args:
            websocket: WebSocket connection
            message: Message content
        """
        try:
            await websocket.send_json(message)
        except Exception:
            # Connection may be closed, remove it
            await self.disconnect(websocket)


async def handle_websocket(
    websocket: WebSocket,
    session_manager: SessionManager,
    ws_manager: WebSocketManager,
    agent_runtime: Any,
    config: Any = None,
) -> None:
    """处理 WebSocket 连接和消息。

    Uses AgentRuntime to process messages — supports tool calls, file ops, etc.
    Integrates three-layer memory system:
      Layer 1: Session Memory (conversation history)
      Layer 2: Workspace prompt files (md → system prompt)
      Layer 3: memorySearch (vector knowledge base)
    """
    import logging
    import os
    from pyclaw.agents.workspace import load_workspace_files, build_workspace_system_prompt
    from pyclaw.memory.mmr import mmr_rerank
    from pyclaw.memory.temporal_decay import apply_temporal_decay

    logger = logging.getLogger(__name__)

    await ws_manager.connect(websocket)

    # --- Layer 1: Session Memory ---
    # Generate a stable session_id for this connection.  The client may pass
    # a previously-issued session_id in the query string so that conversation
    # history is restored after a page refresh (mirrors OpenClaw behaviour).
    query_params = dict(websocket.query_params)
    session_id: str = query_params.get("session_id") or str(uuid.uuid4())

    # Restore conversation history from the JSONL transcript file (if any).
    # This is the key mechanism that makes history survive page refreshes.
    history: list[dict[str, str]] = read_messages(session_id)
    if history:
        logger.info(
            "Restored %d messages from transcript for session %s",
            len(history),
            session_id,
        )

    # Send session_init FIRST so the client can capture the session_id before
    # any other messages arrive.  This mirrors OpenClaw's handshake order.
    await ws_manager.send_to(websocket, {
        "type": "session_init",
        "session_id": session_id,
    })

    # Send restored history to the frontend so it can render previous messages
    if history:
        await ws_manager.send_to(websocket, {
            "type": "session_history",
            "messages": history,
        })

    # Send initial channel status after session_init
    await ws_manager.send_to(websocket, {
        "type": "channel_update",
        "channels": [{"name": "WebSocket", "online": True}],
    })

    # Retrieve memory system components from app state
    app = websocket.app
    memory_manager = getattr(app.state, "memory_manager", None)
    memory_search_config = getattr(app.state, "memory_search_config", None)
    from pyclaw.config.paths import get_paths as _get_paths
    workspace_dir = getattr(app.state, "workspace_dir", str(_get_paths().workspace_dir))

    # ── Per-connection model preference ──────────────────────────────────
    # Use a mutable dict so inner branches can update values without 'nonlocal'.
    # Keys: 'model' (str), 'platform' (str), 'mode' (str)
    conn_model_pref: dict[str, str] = {"model": "", "platform": "", "mode": "standard"}
    
    # Cancellation event for stopping ongoing generation (shared across tasks)
    import asyncio
    cancel_event = asyncio.Event()
    agent_task: asyncio.Task | None = None
    
    async def run_agent_task(
        session_ctx: dict,
        effective_content: str,
        history: list,
        content: str,
    ) -> None:
        """Run agent in a separate task so we can receive cancel messages."""
        nonlocal agent_task
        
        streamed_text = []
        was_cancelled = False
        observability_metadata = None
        
        try:
            async for event in agent_runtime.run(session_ctx, effective_content):
                # Check cancellation
                if cancel_event.is_set():
                    logger.info("Generation cancelled by user")
                    was_cancelled = True
                    await ws_manager.send_to(websocket, {
                        "type": "generation_cancelled",
                        "content": "Generation cancelled",
                    })
                    break
                
                if event.event_type == "text" and event.content:
                    if event.content == "Thinking...":
                        continue
                    if not streamed_text:
                        streamed_text.append(event.content)

                elif event.event_type == "text_delta" and event.content:
                    streamed_text.append(event.content)
                    await ws_manager.send_to(websocket, {
                        "type": "text_delta",
                        "content": event.content,
                    })

                elif event.event_type == "tool_call":
                    await ws_manager.send_to(websocket, {
                        "type": "status",
                        "status": "tool_calling",
                        "content": f"Calling tool: {event.tool_name}",
                    })
                    await ws_manager.send_to(websocket, {
                        "type": "tool_call",
                        "sender": "System",
                        "tool_name": event.tool_name,
                        "tool_args": event.tool_args,
                    })

                elif event.event_type == "tool_result":
                    streamed_text = []
                    await ws_manager.send_to(websocket, {
                        "type": "tool_result",
                        "sender": "System",
                        "tool_name": event.tool_name,
                        "content": event.tool_result or "",
                    })
                    await ws_manager.send_to(websocket, {
                        "type": "status",
                        "status": "thinking",
                        "content": "Continuing to think...",
                    })

                elif event.event_type == "done":
                    if event.metadata:
                        observability_metadata = event.metadata

                elif event.event_type == "error":
                    await ws_manager.send_to(websocket, {
                        "type": "message",
                        "sender": "System",
                        "content": f"Error: {event.content}",
                    })
            
            # Only save to history if not cancelled
            if not was_cancelled:
                final_reply = "".join(streamed_text)
                if final_reply:
                    history.append({"role": "user", "content": content})
                    history.append({"role": "assistant", "content": final_reply})
                    append_message(session_id, "user", content)
                    append_message(session_id, "assistant", final_reply)
                    logger.debug(
                        "Transcript appended for session %s (%d total messages)",
                        session_id, len(history),
                    )
                    
                    if memory_manager and memory_search_config and memory_search_config.enabled:
                        try:
                            await memory_manager.remember(
                                session_id=session_id, agent_name="default",
                                content=content, metadata={"role": "user", "source": "websocket"},
                            )
                            await memory_manager.remember(
                                session_id=session_id, agent_name="default",
                                content=final_reply, metadata={"role": "assistant", "source": "websocket"},
                            )
                            logger.debug("Long-term memory updated for session %s", session_id)
                        except Exception as mem_error:
                            logger.warning("Memory storage failed (non-fatal): %s", mem_error)

                # Append observability panel as a final text_delta before stream_end
                if observability_metadata:
                    from pyclaw.gateway.openai_compat import _format_observability_panel
                    panel_text = _format_observability_panel(observability_metadata)
                    if panel_text:
                        streamed_text.append(panel_text)
                        await ws_manager.send_to(websocket, {
                            "type": "text_delta",
                            "content": panel_text,
                        })

                # 只发送 stream_end 通知前端流式结束，不再重复发送完整内容
                # 因为内容已经通过 text_delta 发送过了
                await ws_manager.send_to(websocket, {
                    "type": "stream_end",
                })
        
        except asyncio.CancelledError:
            logger.info("Agent task was cancelled")
            await ws_manager.send_to(websocket, {
                "type": "generation_cancelled",
                "content": "Generation cancelled",
            })
        finally:
            await ws_manager.send_to(websocket, {
                "type": "generation_ended",
            })
            agent_task = None

    try:
        while True:
            data = await websocket.receive_json()
            logger.info("WS received: %s", data)

            message_type = data.get("type")
            content = data.get("content", "")

            if message_type in ("message", "chat"):
                if not content:
                    continue

                if agent_runtime is not None:
                    # Trim history to avoid context overflow
                    trimmed_history = history[-MAX_HISTORY_MESSAGES:] if len(history) > MAX_HISTORY_MESSAGES else history

                    # --- Layer 2: Build system prompt from workspace files ---
                    workspace_files = load_workspace_files(workspace_dir)
                    workspace_prompt = build_workspace_system_prompt(workspace_files)

                    base_system_prompt = (
                        "You are a helpful AI coding assistant. You have access to the following tools:\n"
                        "- shell: Execute shell commands in the workspace directory\n"
                        "- file_read: Read file contents\n"
                        "- file_write: Write content to files (PLAIN TEXT ONLY)\n"
                        "- web_search: Search the web\n\n"
                        "CRITICAL RULES (you MUST follow these):\n\n"
                        "1. FILE CREATION: When the user asks you to create files, write code, or generate content, "
                        "you MUST use the file_write tool to save it — NEVER output raw code as text.\n\n"
                        "2. CODE EXECUTION: When the user asks you to run code, use the shell tool.\n\n"
                        "3. MULTI-STEP TASKS: You MUST complete ALL steps in a single turn — NEVER stop partway through.\n"
                        "   Example for creating a .pptx file:\n"
                        "   Step 1: shell → pip install python-pptx -q\n"
                        "   Step 2: file_write → save the Python script as a .py file\n"
                        "   Step 3: shell → python3 script.py\n"
                        "   Step 4: Verify the output file was created\n"
                        "   You MUST call ALL required tools in sequence. NEVER stop after just one tool call.\n\n"
                        "4. CONCISE RESPONSES: Keep text responses concise. Do NOT paste large code blocks.\n\n"
                        "5. LANGUAGE: Always respond in the same language the user uses.\n\n"
                        "6. PACKAGE INSTALLATION: If a task requires packages, install them with shell THEN continue with the actual task. "
                        "Installing a package is NEVER the final step — you must proceed to write and execute the code.\n\n"
                        "7. *** BINARY FILE FORMATS (MOST IMPORTANT) ***:\n"
                        "   file_write can ONLY create plain text files (.txt, .py, .md, .json, .html, .css, .js, etc.).\n"
                        "   file_write will REJECT binary formats (.ppt, .pptx, .docx, .xlsx, .pdf, images, etc.).\n"
                        "   To create binary files, you MUST:\n"
                        "   a) Write a Python script using the appropriate library (python-pptx, python-docx, openpyxl)\n"
                        "   b) ALWAYS use `from pptx.util import Pt` for font sizes: `Pt(24)` not `24` or `2400`\n"
                        "   c) Save the script with file_write (as a .py file)\n"
                        "   d) Install the library with shell: `pip install python-pptx -q`\n"
                        "   e) Execute the script with shell: `python3 script.py`\n"
                        "   f) Verify the output file was created successfully\n"
                        "   NEVER try to write binary content directly with file_write — it will fail.\n\n"
                        "   *** PPT GENERATION TEMPLATE (MUST USE THIS PATTERN) ***:\n"
                        "   When creating PPT files, your Python script MUST follow this exact pattern:\n"
                        "   ```\n"
                        "   from pptx import Presentation\n"
                        "   from pptx.util import Pt, Inches\n"
                        "   from pptx.dml.color import RGBColor\n"
                        "   from pptx.enum.text import PP_ALIGN\n"
                        "   prs = Presentation()\n"
                        "   # For each slide:\n"
                        "   slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout\n"
                        "   txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))\n"
                        "   tf = txBox.text_frame\n"
                        "   tf.word_wrap = True\n"
                        "   p = tf.paragraphs[0]\n"
                        "   p.text = 'Your Title'\n"
                        "   p.font.size = Pt(36)\n"
                        "   p.font.bold = True\n"
                        "   # Save the file:\n"
                        "   output_path = '/tmp/output.pptx'\n"
                        "   prs.save(output_path)\n"
                        "   print(f'PPT saved to: {output_path}')\n"
                        "   ```\n"
                        "   CRITICAL: The script MUST actually call prs.save() to create the file.\n"
                        "   NEVER just print a message like 'PPT submitted' without calling prs.save().\n"
                        "   The script must be COMPLETE and SELF-CONTAINED — it must generate the actual .pptx file.\n\n"
                        "8. *** ERROR RECOVERY (MANDATORY — HIGHEST PRIORITY) ***:\n"
                        "   If ANY tool call returns a result containing 'ERROR:', you MUST:\n"
                        "   a) Read and understand the error message carefully\n"
                        "   b) Fix the root cause (not just retry the same thing)\n"
                        "   c) Call the tool again with the corrected parameters\n"
                        "   d) Repeat until the tool succeeds\n"
                        "   e) After fixing a script, you MUST re-execute it with shell to verify the fix works\n"
                        "   NEVER say 'task complete' or give a final answer while there are unresolved errors.\n"
                        "   NEVER skip errors or apologize without fixing them.\n"
                        "   NEVER just rewrite a file without re-running it — the fix is NOT complete until execution succeeds.\n\n"
                        "9. TASK COMPLETION: Only say the task is complete when ALL of these are true:\n"
                        "   a) All tool calls succeeded (no ERROR in any result)\n"
                        "   b) Scripts were executed and produced expected output\n"
                        "   c) Output files were created/verified\n"
                        "   If you wrote or fixed a script, you MUST run it with shell and confirm success before completing."
                    )

                    if workspace_prompt:
                        base_system_prompt = workspace_prompt + "\n\n" + base_system_prompt

                    # --- Layer 3: Search memory knowledge base ---
                    memory_context = ""
                    if (
                        memory_manager
                        and memory_search_config
                        and memory_search_config.enabled
                    ):
                        try:
                            query_config = memory_search_config.query
                            hybrid_config = query_config.hybrid
                            max_results = query_config.max_results
                            min_score = query_config.min_score

                            logger.info(
                                "Memory search starting for query: '%s' (max_results=%d, min_score=%.2f)",
                                content[:80], max_results, min_score,
                            )

                            search_results = await asyncio.wait_for(
                                memory_manager.recall(
                                    query=content,
                                    limit=max_results * 2,
                                ),
                                timeout=5.0,
                            )

                            logger.info(
                                "Memory recall returned %d raw results",
                                len(search_results),
                            )

                            # Convert to dicts for MMR / temporal decay processing
                            result_dicts = [
                                {
                                    "content": r.entry.content,
                                    "score": r.score,
                                    "source": r.source,
                                    "metadata": r.entry.metadata,
                                    "path": r.entry.metadata.get("path", ""),
                                }
                                for r in search_results
                            ]

                            for idx, rd in enumerate(result_dicts[:5]):
                                logger.info(
                                    "  raw[%d] score=%.4f source=%s content=%.60s",
                                    idx, rd["score"], rd["source"], rd["content"],
                                )

                            # Apply MMR reranking if enabled
                            if hybrid_config.mmr.enabled and result_dicts:
                                result_dicts = mmr_rerank(
                                    result_dicts,
                                    lambda_param=hybrid_config.mmr.lambda_,
                                )

                            # Apply temporal decay if enabled
                            if hybrid_config.temporal_decay.enabled and result_dicts:
                                result_dicts = apply_temporal_decay(
                                    result_dicts,
                                    half_life_days=hybrid_config.temporal_decay.half_life_days,
                                )

                            # Filter by min_score and take top K
                            filtered = [
                                r for r in result_dicts
                                if r.get("score", 0) >= min_score
                            ][:max_results]

                            logger.info(
                                "Memory search after filtering: %d results (min_score=%.2f)",
                                len(filtered), min_score,
                            )

                            if filtered:
                                memory_lines = []
                                for idx_m, item in enumerate(filtered):
                                    source_file = item.get("metadata", {}).get("filename", "unknown")
                                    memory_lines.append(
                                        f"{idx_m + 1}. [From {source_file}] {item['content']}"
                                    )
                                memory_context = (
                                    "\n\n<relevant-memories>\n"
                                    "The following memories were recalled from your personal knowledge base. "
                                    "You MUST incorporate this information into your response when it is relevant to the user's query. "
                                    "Refer to these memories naturally as if you already know this about the user. "
                                    "Do NOT ignore them or give generic answers when memories are available.\n"
                                    + "\n".join(memory_lines)
                                    + "\n</relevant-memories>"
                                )
                                logger.info(
                                    "Memory search injected %d results for query: %s",
                                    len(filtered),
                                    content[:80],
                                )
                            else:
                                logger.info("Memory search: no results passed min_score filter")
                        except Exception as mem_error:
                            logger.warning("Memory search failed (non-fatal): %s", mem_error, exc_info=True)

                    # Combine system prompt with memory context
                    full_system_prompt = base_system_prompt
                    if memory_context:
                        full_system_prompt += memory_context

                    # --- Load installed skills into system prompt ---
                    installed_skills = []
                    try:
                        from pyclaw.skills.parser import parse_skill_markdown as _parse_skill
                        from pyclaw.config.paths import get_paths as _get_paths
                        skills_base = _get_paths().workspace_skills_dir
                        if skills_base.exists():
                            for skill_dir in sorted(skills_base.iterdir()):
                                if not skill_dir.is_dir() or skill_dir.name.startswith("."):
                                    continue
                                # 查找 SKILL.md
                                skill_file = None
                                for marker in ["SKILL.md", "skill.md", "skills.md"]:
                                    candidate = skill_dir / marker
                                    if candidate.exists():
                                        skill_file = candidate
                                        break
                                if not skill_file:
                                    continue
                                try:
                                    skill_content = skill_file.read_text(encoding="utf-8")
                                    skill_def = _parse_skill(skill_content, directory_name=skill_dir.name)
                                    if skill_def.name and skill_def.name != "security":
                                        installed_skills.append({
                                            "name": skill_def.name,
                                            "description": skill_def.description,
                                            "content": skill_def.content,
                                            "path": str(skill_dir),
                                        })
                                except Exception as sk_err:
                                    logger.debug("Skip skill %s: %s", skill_dir.name, sk_err)
                        logger.info("Loaded %d installed skills into agent context", len(installed_skills))
                    except Exception as skills_err:
                        logger.warning("Failed to load skills: %s", skills_err)

                    # 将技能内容注入系统提示
                    if installed_skills:
                        from pyclaw.config.paths import get_paths as _get_paths
                        _ws_paths = _get_paths()
                        skills_prompt = f"\n\n## Workspace Info\n"
                        skills_prompt += f"Skills directory: {_ws_paths.workspace_skills_dir}\n"
                        skills_prompt += f"Workspace directory: {_ws_paths.workspace_dir}\n\n"
                        skills_prompt += "## Installed Skills\n"
                        skills_prompt += "You have the following skills installed. Use them when relevant:\n\n"
                        for sk in installed_skills:
                            skills_prompt += f"### Skill: {sk['name']}\n"
                            skills_prompt += f"Installed at: {sk.get('path', 'unknown')}\n"
                            if sk.get('description'):
                                skills_prompt += f"{sk['description']}\n\n"
                            skill_content = sk.get('content', '')
                            if skill_content:
                                # 替换相对路径 scripts/ 为绝对路径
                                skill_path = sk.get('path', '')
                                if skill_path:
                                    skill_content = skill_content.replace(
                                        'scripts/', f"{skill_path}/scripts/"
                                    )
                                # Truncate oversized skill content to avoid exceeding context window
                                if len(skill_content) > MAX_SKILL_CONTENT_LENGTH:
                                    skill_content = skill_content[:MAX_SKILL_CONTENT_LENGTH] + "\n... (truncated)"
                                skills_prompt += f"{skill_content}\n\n"
                        full_system_prompt += skills_prompt

                    # Also prepend memory context to the user message so the LLM
                    # is more likely to incorporate recalled memories into its
                    # response (local models often ignore system-prompt tails).
                    effective_content = content
                    if memory_context:
                        effective_content = (
                            memory_context.strip()
                            + "\n\n---\nUser message:\n"
                            + content
                        )

                    session_ctx = {
                        "history": trimmed_history,
                        "agent_config": {
                            "system_prompt": full_system_prompt,
                        },
                        "skills": installed_skills,
                        "security_policy": {},
                    }
                    # 注入用户在前端选择的模型（ModelSelector.select 会优先读取此字段）
                    if conn_model_pref["model"]:
                        session_ctx["preferred_model"] = conn_model_pref["model"]
                        logger.info(
                            "Using frontend-selected model: platform=%s, model=%s",
                            conn_model_pref["platform"], conn_model_pref["model"],
                        )
                    
                    # Adjust LLM parameters based on mode
                    mode = conn_model_pref.get("mode", "standard")
                    mode_params = {
                        "standard": {"temperature": DEFAULT_TEMPERATURE, "max_tokens": DEFAULT_MAX_TOKENS},
                        "fast": {"temperature": 0.8, "max_tokens": 1024},
                        "precise": {"temperature": 0.3, "max_tokens": 4096},
                    }
                    session_ctx["mode_params"] = mode_params.get(mode, mode_params["standard"])
                    logger.info("Using mode: %s, params: %s", mode, session_ctx["mode_params"])

                    # Send thinking status before starting
                    await ws_manager.send_to(websocket, {
                        "type": "status",
                        "status": "thinking",
                        "content": "Thinking...",
                    })
                    
                    # Notify frontend that generation has started
                    await ws_manager.send_to(websocket, {
                        "type": "generation_started",
                    })

                    # Reset cancel event and start agent as background task
                    cancel_event.clear()
                    agent_task = asyncio.create_task(
                        run_agent_task(session_ctx, effective_content, history, content)
                    )
                    # Don't await here - continue receiving messages
                else:
                    await ws_manager.send_to(websocket, {
                        "type": "message",
                        "sender": "System",
                        "content": "Agent runtime not initialized.",
                    })

            elif message_type == "set_model":
                # 前端切换模型/平台
                new_platform = data.get("platform", "")
                new_model = data.get("model", "")
                logger.info("WS set_model: platform=%s, model=%s", new_platform, new_model)

                # 检查是否真正变化
                old_platform = conn_model_pref.get("platform", "")
                old_model = conn_model_pref.get("model", "")
                model_changed = (new_model and new_model != old_model) or (new_platform and new_platform != old_platform)

                # 更新 per-connection 可变 dict，后续请求的 session_ctx 会注入 preferred_model
                if new_model:
                    conn_model_pref["model"] = new_model
                if new_platform:
                    conn_model_pref["platform"] = new_platform

                # 动态注册模型到 model_selector，使 ModelSelector._get_model_config 能找到它
                if agent_runtime is not None and new_model:
                    model_selector = getattr(agent_runtime, "model_selector", None)
                    if model_selector is not None:
                        models_dict = model_selector.models_config.setdefault("models", {})
                        # 根据平台决定 base_url 和 provider
                        _platform = new_platform or "internal"
                        # Look up provider from unified config
                        app_config = getattr(websocket.app.state, 'config', None)
                        provider_entry = None
                        if app_config and hasattr(app_config, 'llm') and hasattr(app_config.llm, 'providers'):
                            provider_entry = app_config.llm.providers.get(_platform)

                        if provider_entry:
                            _base_url = provider_entry.base_url
                            _api_key = provider_entry.api_key
                            _provider = _platform
                            _provider_type = provider_entry.type
                            # Find model-specific settings
                            model_max_tokens = None
                            model_temperature = None
                            for m in provider_entry.models:
                                if m.id == new_model:
                                    model_max_tokens = m.max_tokens
                                    model_temperature = m.temperature
                                    break
                        else:
                            # Fallback for local/ollama
                            from pyclaw.constants import DEFAULT_LOCAL_LLM_BASE_URL
                            _base_url = app_config.llm.get_local_base_url() if app_config else DEFAULT_LOCAL_LLM_BASE_URL
                            _api_key = None
                            _provider = 'local'
                            _provider_type = 'ollama'
                            model_max_tokens = None
                            model_temperature = None

                        if new_model not in models_dict:
                            models_dict[new_model] = {
                                "provider": _provider,
                                "provider_type": _provider_type,
                                "base_url": _base_url,
                                "api_key": _api_key,
                                "max_tokens": model_max_tokens or 16384,
                                "temperature": model_temperature or 0.7,
                            }
                            logger.info(
                                "Dynamically registered model: %s (provider=%s, provider_type=%s, base_url=%s)",
                                new_model, _provider, _provider_type, _base_url,
                            )
                        else:
                            # 已存在，更新 provider、provider_type 和 base_url
                            models_dict[new_model]["provider"] = _provider
                            models_dict[new_model]["provider_type"] = _provider_type
                            models_dict[new_model]["base_url"] = _base_url
                            if _api_key is not None:
                                models_dict[new_model]["api_key"] = _api_key
                        # 清除缓存，让 select 重新加载
                        model_selector._models_cache.pop(new_model, None)
                        logger.info(
                            "Model selector updated: preferred_model=%s, platform=%s",
                            conn_model_pref["model"], conn_model_pref["platform"],
                        )

                # Sync to global model preference so HTTP API (DingTalk etc.) picks it up
                if model_changed:
                    global_pref = getattr(websocket.app.state, "global_model_pref", None)
                    if global_pref is not None:
                        if new_model:
                            global_pref["model"] = new_model
                        if new_platform:
                            global_pref["platform"] = new_platform
                        logger.info(
                            "Global model pref updated: model=%s, platform=%s",
                            global_pref["model"], global_pref["platform"],
                        )

                await ws_manager.send_to(websocket, {
                    "type": "model_updated",
                    "platform": new_platform,
                    "model": new_model,
                    "message": f"Switched to {new_platform.upper()} / {new_model}" if model_changed else "",
                    "changed": model_changed,
                })

            elif message_type == "set_mode":
                # Frontend switches mode (standard/fast/precise)
                new_mode = data.get("mode", "standard")
                conn_model_pref["mode"] = new_mode
                mode_labels = {
                    "standard": "Standard mode",
                    "fast": "Fast mode",
                    "precise": "Precise mode",
                }
                logger.info("WS set_mode: %s", new_mode)
                await ws_manager.send_to(websocket, {
                    "type": "mode_updated",
                    "mode": new_mode,
                    "message": f"Switched to {mode_labels.get(new_mode, new_mode)}",
                })

            elif message_type == "get_skills":
                # 前端请求已安装技能列表
                # 直接扫描 skills 目录下的 .md 文件，不走 SkillLoader.load_all()
                # 避免 _validate_security_skill 在 security 技能不存在时抛出 RuntimeError
                skill_list = []
                try:
                    from pyclaw.skills.parser import parse_skill_markdown
                    from pyclaw.config.paths import get_paths as _get_paths
                    skills_dir = _get_paths().skills_dir
                    if skills_dir.exists():
                        for skill_file in sorted(skills_dir.glob("*.md")):
                            try:
                                skill = parse_skill_markdown(skill_file.read_text(encoding="utf-8"))
                                if skill.name and skill.name != "security":
                                    skill_list.append({
                                        "name": skill.name,
                                        "description": getattr(skill, "description", ""),
                                        "enabled": True,
                                    })
                            except Exception as parse_err:
                                logger.debug("Skipping skill file %s: %s", skill_file.name, parse_err)
                    logger.info("get_skills: found %d user skills in %s", len(skill_list), skills_dir)
                except Exception as skills_err:
                    logger.warning("get_skills failed: %s", skills_err)

                await ws_manager.send_to(websocket, {
                    "type": "skills_list",
                    "skills": skill_list,
                })

            elif message_type == "set_debug":
                # 前端一键开关 debug 日志
                enabled: bool = bool(data.get("enabled", False))
                logger.info("WS set_debug: enabled=%s", enabled)

                # 切换 internal_client 的 debug 开关
                try:
                    from pyclaw.llm.internal_client import set_debug_logging
                    set_debug_logging(enabled)
                except ImportError:
                    pass

                # 同时调整根 logger 级别
                import logging as _logging
                root_logger = _logging.getLogger("pyclaw")
                root_logger.setLevel(_logging.DEBUG if enabled else _logging.INFO)

                await ws_manager.send_to(websocket, {
                    "type": "debug_updated",
                    "enabled": enabled,
                    "message": f"Debug logging {'enabled' if enabled else 'disabled'}",
                })

            elif message_type == "ping":
                await ws_manager.send_to(websocket, {"type": "pong"})

            elif message_type == "cancel":
                # User requested to cancel ongoing generation
                cancel_event.set()
                logger.info("Cancel requested by user")
                
                # Also cancel the task if it's running
                if agent_task and not agent_task.done():
                    agent_task.cancel()
                    logger.info("Agent task cancelled")
                
                await ws_manager.send_to(websocket, {
                    "type": "cancel_acknowledged",
                    "message": "Cancelling...",
                })

    except Exception as exc:
        # WebSocketDisconnect is normal — don't log as error
        from starlette.websockets import WebSocketDisconnect
        if isinstance(exc, WebSocketDisconnect):
            logger.info("WebSocket disconnected: code=%s", exc.code)
        else:
            logger.error("WebSocket handler error: %s", exc, exc_info=True)
    finally:
        await ws_manager.disconnect(websocket)
