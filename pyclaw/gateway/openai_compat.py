"""
OpenAI-compatible Chat Completions API

Implements /v1/chat/completions endpoint, compatible with OpenAI API format.
Used for integration with DingTalk, WeChat, Feishu, and other channels.
"""

import asyncio
import json
import logging
import os
import time
import uuid
from typing import Any, AsyncGenerator, Dict, List, Optional, Union

from fastapi import APIRouter, Depends, HTTPException, Header, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field

from pyclaw.constants import (
    AGENT_TIMEOUT_SECONDS,
    DEFAULT_MAX_TOKENS,
    DEFAULT_TEMPERATURE,
    DINGTALK_CLIENT_TIMEOUT,
    ERROR_MESSAGE_PREFIXES,
    MAX_HISTORY_MESSAGES,
)

logger = logging.getLogger(__name__)


# ─── Request/Response Models ────────────────────────────────────────────────


class ChatMessage(BaseModel):
    """Chat message"""
    role: str = Field(..., description="Role: system, user, assistant")
    content: str = Field(..., description="Message content")
    name: Optional[str] = Field(None, description="Sender name")


class ChatCompletionRequest(BaseModel):
    """Chat completion request"""
    model: str = Field(default="default", description="Model name")
    messages: List[ChatMessage] = Field(..., description="Message list")
    temperature: Optional[float] = Field(default=0.7, ge=0, le=2)
    max_tokens: Optional[int] = Field(default=DEFAULT_MAX_TOKENS, ge=1)
    stream: bool = Field(default=False, description="Whether to stream output")
    user: Optional[str] = Field(None, description="User identifier")


class ChatCompletionChoice(BaseModel):
    """Completion choice"""
    index: int = 0
    message: ChatMessage
    finish_reason: str = "stop"


class ChatCompletionUsage(BaseModel):
    """Token usage"""
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0


class ChatCompletionResponse(BaseModel):
    """Chat completion response"""
    id: str = Field(default_factory=lambda: f"chatcmpl-{uuid.uuid4().hex[:12]}")
    object: str = "chat.completion"
    created: int = Field(default_factory=lambda: int(time.time()))
    model: str = "pyclaw"
    choices: List[ChatCompletionChoice]
    usage: ChatCompletionUsage = Field(default_factory=ChatCompletionUsage)


class ChatCompletionChunkDelta(BaseModel):
    """Streaming response delta"""
    role: Optional[str] = None
    content: Optional[str] = None


class ChatCompletionChunkChoice(BaseModel):
    """Streaming response choice"""
    index: int = 0
    delta: ChatCompletionChunkDelta
    finish_reason: Optional[str] = None


class ChatCompletionChunk(BaseModel):
    """Streaming response chunk"""
    id: str
    object: str = "chat.completion.chunk"
    created: int
    model: str = "pyclaw"
    choices: List[ChatCompletionChunkChoice]


# ─── Auth ───────────────────────────────────────────────────────────────────


def verify_auth_token(
    authorization: Optional[str] = Header(None),
    x_openclaw_token: Optional[str] = Header(None, alias="X-OpenClaw-Token"),
    x_pyclaw_token: Optional[str] = Header(None, alias="X-PyClaw-Token"),
) -> Optional[str]:
    """Verify authentication token
    
    Supports multiple authentication methods:
    - Authorization: Bearer <token>
    - X-OpenClaw-Token: <token>
    - X-PyClaw-Token: <token>
    
    Returns:
        Extracted token, or None if not provided
    """
    if authorization and authorization.startswith("Bearer "):
        return authorization[7:]
    if x_pyclaw_token:
        return x_pyclaw_token
    if x_openclaw_token:
        return x_openclaw_token
    return None


# ─── Router Factory ─────────────────────────────────────────────────────────


def create_chat_completions_router(
    config: Any,
    agent_runtime: Any,
) -> APIRouter:
    """Create chat completions router
    
    Args:
        config: PyClaw configuration
        agent_runtime: Agent runtime
        
    Returns:
        FastAPI router
    """
    router = APIRouter(prefix="/v1", tags=["chat"])
    
    # 从配置获取认证 token
    gateway_auth_token: Optional[str] = None
    if hasattr(config, 'gateway') and hasattr(config.gateway, 'auth'):
        auth = config.gateway.auth
        if hasattr(auth, 'token'):
            token = auth.token
            gateway_auth_token = token.get_secret_value() if hasattr(token, 'get_secret_value') else str(token)
    
    async def check_auth(token: Optional[str] = Depends(verify_auth_token)) -> None:
        """Check authentication"""
        if gateway_auth_token:
            if not token:
                raise HTTPException(status_code=401, detail="Missing authentication token")
            if token != gateway_auth_token:
                raise HTTPException(status_code=401, detail="Invalid authentication token")
    
    @router.post("/chat/completions", response_model=None)
    async def chat_completions(
        request: ChatCompletionRequest,
        fastapi_request: Request,
        _auth: None = Depends(check_auth),
    ) -> Union[ChatCompletionResponse, StreamingResponse]:
        """OpenAI-compatible chat completion API
        
        Supports both streaming and non-streaming responses.
        """
        logger.info("Chat completions request: model=%s, stream=%s", request.model, request.stream)
        
        # 获取 memory_manager （用于注入长期记忆）
        memory_manager = getattr(fastapi_request.app.state, 'memory_manager', None)
        
        # 获取 cron_scheduler （用于定时任务工具）
        cron_scheduler = getattr(fastapi_request.app.state, 'cron_scheduler', None)
        
        # 动态注册本地 Ollama 模型（确保 model_selector 能找到它）
        if agent_runtime is not None:
            model_selector = getattr(agent_runtime, "model_selector", None)
            if model_selector is not None:
                models_dict = model_selector.models_config.setdefault("models", {})
                local_model = os.environ.get("PYCLAW_DEFAULT_LOCAL_MODEL", "qwen3.5:9B")
                if local_model not in models_dict:
                    from pyclaw.config.loader import load_config as _load_config
                    models_dict[local_model] = {
                        "provider": "local",
                        "base_url": _load_config().llm.get_local_base_url(),
                        "api_key": None,
                        "max_tokens": 4096,
                        "temperature": 0.7,
                    }
                    logger.info("Dynamically registered local model: %s", local_model)
                # 清除缓存
                if hasattr(model_selector, '_models_cache'):
                    model_selector._models_cache.pop(local_model, None)
        
        # 提取用户消息
        user_message = ""
        for msg in reversed(request.messages):
            if msg.role == "user":
                user_message = msg.content
                break
        
        if not user_message:
            raise HTTPException(status_code=400, detail="No user message found")
        
        # 构建历史
        history = []
        for msg in request.messages:
            history.append({"role": msg.role, "content": msg.content})
        
        # 构建完整的会话上下文（和 WebSocket 相同的路径）
        session_ctx = await _build_full_session_ctx(
            session_id=request.user or f"api-{uuid.uuid4().hex[:8]}",
            history=history[:-1],  # 不包括最后一条用户消息
            user_message=user_message,
            request=request,
            memory_manager=memory_manager,  # 传入记忆管理器
            cron_scheduler=cron_scheduler,  # 传入定时任务调度器
            app_state=fastapi_request.app.state,  # 传入 app state 用于读取全局模型偏好
        )
        
        if request.stream:
            return StreamingResponse(
                _stream_response(agent_runtime, session_ctx, user_message, request.model),
                media_type="text/event-stream",
            )
        else:
            return await _non_stream_response(agent_runtime, session_ctx, user_message, request.model)
    
    @router.get("/models")
    async def list_models():
        """List available models"""
        return {
            "object": "list",
            "data": [
                {
                    "id": "pyclaw",
                    "object": "model",
                    "created": int(time.time()),
                    "owned_by": "pyclaw",
                },
            ],
        }
    
    return router


def _filter_error_history(history: list) -> list:
    """Remove assistant messages that contain error/timeout content.

    These pollute the conversation context and cause the LLM to
    reproduce error messages instead of calling tools.

    Args:
        history: List of message dicts with 'role' and 'content'.

    Returns:
        Filtered list with error assistant messages (and their
        preceding user messages) removed.
    """
    indices_to_remove: set[int] = set()
    for idx, msg in enumerate(history):
        if msg.get("role") != "assistant":
            continue
        content = msg.get("content", "")
        if any(content.startswith(prefix) for prefix in ERROR_MESSAGE_PREFIXES):
            indices_to_remove.add(idx)
            # Also remove the preceding user message that triggered
            # the error, so the LLM doesn't see a dangling question
            # without a proper answer.
            if idx > 0 and history[idx - 1].get("role") == "user":
                indices_to_remove.add(idx - 1)

    if indices_to_remove:
        logger.info("Filtered %d error messages from history", len(indices_to_remove))

    return [msg for idx, msg in enumerate(history) if idx not in indices_to_remove]


# ─── 构建完整会话上下文（和 WebSocket 相同） ─────────────────────────────────

async def _build_full_session_ctx(
    session_id: str,
    history: list,
    user_message: str,
    request: ChatCompletionRequest,
    memory_manager: Any = None,
    cron_scheduler: Any = None,
    app_state: Any = None,
) -> dict:
    """构建完整的会话上下文，和 WebSocket 走完全相同的路径
    
    包含：
    - 完整的系统提示（工具说明、规则等）
    - 已安装的技能
    - 长期记忆（从 MemoryManager 检索）
    - 会话历史
    - 模式参数
    """
    import os
    from pathlib import Path
    
    # 基础系统提示（和 websocket.py 完全相同）
    base_system_prompt = (
        "You are a helpful AI coding assistant. You have access to the following tools:\n"
        "- shell: Execute shell commands in the workspace directory\n"
        "- file_read: Read file contents\n"
        "- file_write: Write content to files (PLAIN TEXT ONLY)\n"
        "- web_search: Search the web\n"
        "- memory_save: Save important information to long-term memory\n"
        "- memory_search: Search long-term memory for past information\n\n"
        "CRITICAL RULES (you MUST follow these):\n\n"
        "1. FILE CREATION: When the user asks you to create files, write code, or generate content, "
        "you MUST use the file_write tool to save it — NEVER output raw code as text.\n\n"
        "2. CODE EXECUTION: When the user asks you to run code, use the shell tool.\n\n"
        "3. MEMORY: When the user tells you something important to remember (like their name, "
        "preferences, or key facts), use the memory_save tool IMMEDIATELY. "
        "When you need to recall past information, use memory_search first.\n\n"
        "4. MULTI-STEP TASKS: You MUST complete ALL steps in a single turn — NEVER stop partway through.\n"
        "   Example for creating a .pptx file:\n"
        "   Step 1: shell → pip install python-pptx -q\n"
        "   Step 2: file_write → save the Python script as a .py file\n"
        "   Step 3: shell → python3 script.py\n"
        "   Step 4: Verify the output file was created\n"
        "   You MUST call ALL required tools in sequence. NEVER stop after just one tool call.\n\n"
        "5. CONCISE RESPONSES: Keep text responses concise. Do NOT paste large code blocks.\n\n"
        "6. LANGUAGE: Always respond in the same language the user uses.\n\n"
        "7. PACKAGE INSTALLATION: If a task requires packages, install them with shell THEN continue with the actual task. "
        "Installing a package is NEVER the final step — you must proceed to write and execute the code.\n\n"
        "8. *** BINARY FILE FORMATS (MOST IMPORTANT) ***:\n"
        "   file_write can ONLY create plain text files (.txt, .py, .md, .json, .html, .css, .js, etc.).\n"
        "   file_write will REJECT binary formats (.ppt, .pptx, .docx, .xlsx, .pdf, images, etc.).\n"
        "   To create binary files, you MUST:\n"
        "   a) Write a Python script using the appropriate library (python-pptx, python-docx, openpyxl)\n"
        "   b) ALWAYS use `from pptx.util import Pt` for font sizes: `Pt(24)` not `24` or `2400`\n"
        "   c) Save the script with file_write (as a .py file)\n"
        "   d) Install the library with shell: `pip install python-pptx -q`\n"
        "   e) Execute the script with shell: `python3 script.py`\n"
        "   f) Verify the output file was created successfully\n"
        "   NEVER try to write binary content directly with file_write — it will fail.\n\n"
        "   *** PPT GENERATION TEMPLATE (MUST USE THIS PATTERN) ***:\n"
        "   When creating PPT files, your Python script MUST follow this exact pattern:\n"
        "   ```\n"
        "   from pptx import Presentation\n"
        "   from pptx.util import Pt, Inches\n"
        "   from pptx.dml.color import RGBColor\n"
        "   from pptx.enum.text import PP_ALIGN\n"
        "   prs = Presentation()\n"
        "   # For each slide:\n"
        "   slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout\n"
        "   txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))\n"
        "   tf = txBox.text_frame\n"
        "   tf.word_wrap = True\n"
        "   p = tf.paragraphs[0]\n"
        "   p.text = 'Your Title'\n"
        "   p.font.size = Pt(36)\n"
        "   p.font.bold = True\n"
        "   # Save the file:\n"
        "   output_path = '/tmp/output.pptx'\n"
        "   prs.save(output_path)\n"
        "   print(f'PPT saved to: {output_path}')\n"
        "   ```\n"
        "   CRITICAL: The script MUST actually call prs.save() to create the file.\n"
        "   NEVER just print a message like 'PPT submitted' without calling prs.save().\n"
        "   The script must be COMPLETE and SELF-CONTAINED — it must generate the actual .pptx file.\n\n"
        "9. *** ERROR RECOVERY (MANDATORY — HIGHEST PRIORITY) ***:\n"
        "   If ANY tool call returns a result containing 'ERROR:', you MUST:\n"
        "   a) Read and understand the error message carefully\n"
        "   b) Fix the root cause (not just retry the same thing)\n"
        "   c) Call the tool again with the corrected parameters\n"
        "   d) Repeat until the tool succeeds\n"
        "   e) After fixing a script, you MUST re-execute it with shell to verify the fix works\n"
        "   NEVER say 'task complete' or give a final answer while there are unresolved errors.\n"
        "   NEVER skip errors or apologize without fixing them.\n"
        "   NEVER just rewrite a file without re-running it — the fix is NOT complete until execution succeeds.\n\n"
        "10. TASK COMPLETION: Only say the task is complete when ALL of these are true:\n"
        "   a) All tool calls succeeded (no ERROR in any result)\n"
        "   b) Scripts were executed and produced expected output\n"
        "   c) Output files were created/verified\n"
        "   If you wrote or fixed a script, you MUST run it with shell and confirm success before completing.\n\n"
        "11. WORKSPACE PATHS:\n"
        "   - Skills directory: Use the path shown in 'Workspace Info' section below\n"
        "   - To create a new skill: write SKILL.md to {skills_dir}/{skill-name}/SKILL.md\n"
        "   - All file operations are restricted to the workspace directory\n"
    )
    
    # 加载已安装的技能
    installed_skills = []
    try:
        from pyclaw.skills.parser import parse_skill_markdown as _parse_skill
        from pyclaw.config.paths import get_paths as _get_paths
        skills_base = _get_paths().workspace_skills_dir
        if skills_base.exists():
            for skill_dir in sorted(skills_base.iterdir()):
                if not skill_dir.is_dir() or skill_dir.name.startswith("."):
                    continue
                # 查找 SKILL.md
                skill_file = None
                for marker in ["SKILL.md", "skill.md", "skills.md"]:
                    candidate = skill_dir / marker
                    if candidate.exists():
                        skill_file = candidate
                        break
                if not skill_file:
                    continue
                try:
                    skill_content = skill_file.read_text(encoding="utf-8")
                    skill_def = _parse_skill(skill_content, directory_name=skill_dir.name)
                    if skill_def.name and skill_def.name != "security":
                        installed_skills.append({
                            "name": skill_def.name,
                            "description": skill_def.description,
                            "content": skill_def.content,
                            "path": str(skill_dir),
                        })
                except Exception as sk_err:
                    logger.debug("Skip skill %s: %s", skill_dir.name, sk_err)
        logger.info("Loaded %d installed skills for API request", len(installed_skills))
    except Exception as skills_err:
        logger.warning("Failed to load skills: %s", skills_err)
    
    # 将技能内容注入系统提示
    full_system_prompt = base_system_prompt
    if installed_skills:
        from pyclaw.config.paths import get_paths as _get_paths
        _ws_paths = _get_paths()
        skills_prompt = f"\n\n## Workspace Info\n"
        skills_prompt += f"Skills directory: {_ws_paths.workspace_skills_dir}\n"
        skills_prompt += f"Workspace directory: {_ws_paths.workspace_dir}\n\n"
        skills_prompt += "## Installed Skills\n"
        skills_prompt += (
            "You have the following skills installed as function tools. "
            "IMPORTANT: When a user asks about topics covered by these skills "
            "(e.g. stock analysis, weather, web search), you MUST call the "
            "corresponding skill tool (e.g. skill_a-stock-analyzer, skill_buffett) "
            "to get real-time data. Do NOT call the shell tool directly to run "
            "skill scripts — always use the skill_* function tools which handle "
            "script execution, orchestration chains, and data passing automatically. "
            "Do NOT answer from memory or prior conversation history.\n\n"
            "ORCHESTRATION RULES:\n"
            "- Skills with chain_after define an automatic follow-up chain.\n"
            "- When a provider skill (e.g. a-stock-analyzer) has chain_after: [buffett], "
            "you MUST call skill_buffett AFTER getting data from skill_a-stock-analyzer, "
            "passing the data output as the 'context' parameter.\n"
            "- Consumer skills (role=consumer) require a 'context' parameter with upstream data.\n\n"
        )
        # Import orchestration registry to get trigger_description and chain_after
        from pyclaw.skills.tool_adapter import get_skill_orchestration_registry
        orch_registry = get_skill_orchestration_registry()

        for sk in installed_skills:
            tool_name = f"skill_{sk['name']}"
            orch = orch_registry.get(tool_name, {})
            trigger = orch.get("trigger_description", "")
            chain_after = orch.get("chain_after", [])
            role = orch.get("role", "")

            line = f"- **{tool_name}**: {sk.get('description', 'No description')}"
            if trigger:
                line += f" | WHEN TO CALL: {trigger}"
            if chain_after:
                chain_tools = ", ".join(f"skill_{c}" for c in chain_after)
                line += f" | THEN CALL: {chain_tools}"
            if role == "consumer":
                line += " | REQUIRES: context parameter with upstream data"
            skills_prompt += line + "\n"
        full_system_prompt += skills_prompt
    
    # === 长期记忆注入（借鉴 OpenClaw/Claude-Code） ===
    # 从 MemoryManager 检索与用户或当前消息相关的记忆
    if memory_manager is not None:
        try:
            # 检索与用户消息相关的记忆（语义搜索），5秒超时避免阻塞
            memories = await asyncio.wait_for(
                memory_manager.recall(
                    query=user_message,
                    session_id=session_id,
                    limit=5,
                ),
                timeout=5.0,
            )
            
            if memories:
                memory_prompt = "\n\n## Long-term Memory\n"
                memory_prompt += "You have the following memories about this user or previous conversations:\n\n"
                for mem in memories:
                    # 过滤低相关性的记忆
                    if mem.score < 0.3:
                        continue
                    memory_prompt += f"- {mem.entry.content}\n"
                full_system_prompt += memory_prompt
                logger.info("Injected %d long-term memories into system prompt", len(memories))
        except asyncio.TimeoutError:
            logger.warning("Memory retrieval timed out (5s), skipping memory injection")
        except Exception as mem_err:
            logger.warning("Memory retrieval failed: %s: %s", type(mem_err).__name__, mem_err)
    
    # Get model from global runtime preference (set by WebUI via WebSocket),
    # falling back to static config if not available.
    global_pref = getattr(app_state, "global_model_pref", None) if app_state else None
    if global_pref and global_pref.get("model"):
        default_model = global_pref["model"]
        default_platform = global_pref.get("platform", "local")
        logger.info("Using global model pref: model=%s, platform=%s", default_model, default_platform)
    else:
        default_model = "qwen3.5:9B"
        default_platform = "local"
        try:
            from pyclaw.config.loader import load_config
            config = load_config()
            if config and hasattr(config, 'llm') and config.llm.providers:
                default_prov = config.llm.providers.get(config.llm.default)
                if default_prov and default_prov.models:
                    default_model = default_prov.models[0].id
                default_platform = config.llm.default
        except Exception as config_err:
            logger.debug("Failed to load config for default model: %s", config_err)

    # 构建 session_ctx
    session_ctx = {
        "session_id": session_id,
        "history": _filter_error_history(history[-MAX_HISTORY_MESSAGES:] if len(history) > MAX_HISTORY_MESSAGES else history),
        "agent_config": {
            "system_prompt": full_system_prompt,
        },
        "skills": installed_skills,
        "security_policy": {},
        "mode_params": {
            "temperature": request.temperature or DEFAULT_TEMPERATURE,
            "max_tokens": request.max_tokens or DEFAULT_MAX_TOKENS,
        },
        # 默认使用本地 Ollama 模型
        "preferred_model": default_model,
        "preferred_platform": default_platform,
        # 传入 memory_manager 和 cron_scheduler 供工具使用
        "memory_manager": memory_manager,
        "cron_scheduler": cron_scheduler,
    }
    
    # 如果请求指定了其他模型，则覆盖
    if request.model and request.model != "default":
        session_ctx["preferred_model"] = request.model
    
    return session_ctx


async def _stream_response(
    agent_runtime: Any,
    session_ctx: dict,
    message: str,
    model: str,
) -> AsyncGenerator[str, None]:
    """生成流式响应
    
    Args:
        agent_runtime: Agent 运行时
        session_ctx: 会话上下文
        message: 用户消息
        model: 模型名称
        
    Yields:
        SSE 格式的响应块
    """
    response_id = f"chatcmpl-{uuid.uuid4().hex[:12]}"
    created = int(time.time())
    
    # 发送角色
    chunk = ChatCompletionChunk(
        id=response_id,
        created=created,
        model=model,
        choices=[
            ChatCompletionChunkChoice(
                delta=ChatCompletionChunkDelta(role="assistant"),
            )
        ],
    )
    yield f"data: {chunk.model_dump_json()}\n\n"
    
    full_content = ""
    
    try:
        async for event in agent_runtime.run(session_ctx, message):
            if event.event_type == "text" and event.content:
                full_content += event.content
                chunk = ChatCompletionChunk(
                    id=response_id,
                    created=created,
                    model=model,
                    choices=[
                        ChatCompletionChunkChoice(
                            delta=ChatCompletionChunkDelta(content=event.content),
                        )
                    ],
                )
                yield f"data: {chunk.model_dump_json()}\n\n"
    except Exception as e:
        logger.error("Stream error: %s", e, exc_info=True)
        error_chunk = ChatCompletionChunk(
            id=response_id,
            created=created,
            model=model,
            choices=[
                ChatCompletionChunkChoice(
                    delta=ChatCompletionChunkDelta(content=f"\n\n[Error: {str(e)}]"),
                )
            ],
        )
        yield f"data: {error_chunk.model_dump_json()}\n\n"
    
    # 发送结束标记
    final_chunk = ChatCompletionChunk(
        id=response_id,
        created=created,
        model=model,
        choices=[
            ChatCompletionChunkChoice(
                delta=ChatCompletionChunkDelta(),
                finish_reason="stop",
            )
        ],
    )
    yield f"data: {final_chunk.model_dump_json()}\n\n"
    yield "data: [DONE]\n\n"


def _format_markdown_table(content: str) -> str:
    """优化 Markdown 表格格式，确保对齐和可读性。

    Args:
        content: 包含 Markdown 表格的文本内容

    Returns:
        格式化后的文本，表格列对齐、分隔符统一
    """
    import re
    
    lines = content.split('\n')
    formatted_lines = []
    i = 0
    
    while i < len(lines):
        line = lines[i]
        
        # 检测表格起始行（包含 | 且至少有 2 个 |）
        if '|' in line and line.count('|') >= 2:
            # 收集整个表格
            table_lines = [line]
            j = i + 1
            
            # 继续收集表格行（包含 | 的行或分隔符行如 |---|---|）
            while j < len(lines):
                next_line = lines[j]
                if '|' in next_line or (next_line.strip() and re.match(r'^[\s|-]+$', next_line)):
                    table_lines.append(next_line)
                    j += 1
                else:
                    break
            
            # 如果表格行数 >= 2（至少一行表头 + 一行数据），则进行格式化
            if len(table_lines) >= 2:
                try:
                    formatted_table = _parse_and_format_table(table_lines)
                    formatted_lines.extend(formatted_table)
                    i = j
                    continue
                except Exception as e:
                    logger.debug("Table formatting failed: %s, using original", e)
            
            # 格式化失败或不是有效表格，保留原样
            formatted_lines.append(line)
            i += 1
        else:
            formatted_lines.append(line)
            i += 1
    
    return '\n'.join(formatted_lines)


def _parse_and_format_table(table_lines: List[str]) -> List[str]:
    """解析并格式化单个 Markdown 表格。

    Args:
        table_lines: 表格的原始行列表

    Returns:
        格式化后的表格行列表
    """
    import re
    
    # 解析每一行的单元格
    rows = []
    for line in table_lines:
        # 移除行首尾的 | 并分割
        cells = [cell.strip() for cell in line.strip().split('|')]
        # 移除空的首尾单元格（由 | 开头和结尾导致）
        if cells and cells[0] == '':
            cells = cells[1:]
        if cells and cells[-1] == '':
            cells = cells[:-1]
        if cells:  # 非空行
            rows.append(cells)
    
    if not rows:
        return table_lines
    
    # 计算每列的最大宽度
    num_cols = max(len(row) for row in rows)
    col_widths = [0] * num_cols
    
    for row in rows:
        for idx, cell in enumerate(row):
            # 计算单元格宽度（中文字符按 2 个英文字符宽度计算）
            width = 0
            for char in cell:
                if '\u4e00' <= char <= '\u9fff':  # 中文字符
                    width += 2
                else:
                    width += 1
            col_widths[idx] = max(col_widths[idx], width)
    
    # 格式化每一行
    formatted_rows = []
    for row_idx, row in enumerate(rows):
        # 补齐缺失的列
        while len(row) < num_cols:
            row.append('')
        
        formatted_cells = []
        for idx, cell in enumerate(row):
            # 计算填充宽度（考虑中英文混合）
            cell_width = 0
            for char in cell:
                if '\u4e00' <= char <= '\u9fff':
                    cell_width += 2
                else:
                    cell_width += 1
            
            padding = col_widths[idx] - cell_width
            # 表头分隔符行使用 - 填充，其他行右对齐
            if row_idx == 1 and re.match(r'^-+:?-*$', cell):
                # 分隔符行，保持原有对齐方式
                if cell.startswith(':') and cell.endswith(':'):
                    formatted_cell = ':' + '-' * (col_widths[idx] - 2) + ':'
                elif cell.startswith(':'):
                    formatted_cell = ':' + '-' * (col_widths[idx] - 1)
                elif cell.endswith(':'):
                    formatted_cell = '-' * (col_widths[idx] - 1) + ':'
                else:
                    formatted_cell = '-' * col_widths[idx]
            else:
                # 数据行右对齐
                formatted_cell = ' ' * padding + cell
            formatted_cells.append(formatted_cell)
        
        formatted_rows.append('| ' + ' | '.join(formatted_cells) + ' |')
    
    return formatted_rows


def _format_observability_panel(data: dict) -> str:
    """Format observability metrics into a Markdown panel.

    Args:
        data: Observability dict from AgentEvent.metadata

    Returns:
        Formatted Markdown string to append to the response.
    """
    session_id = data.get("session_id", "unknown")
    tool_chain = data.get("tool_chain", [])
    tool_calls_count = data.get("tool_calls_count", 0)
    total_tool_ms = data.get("total_tool_ms", 0)
    total_llm_ms = data.get("total_llm_ms", 0)
    session_ms = data.get("session_ms", 0)
    tool_freq = data.get("tool_frequency", {})
    slowest_name = data.get("slowest_tool_name", "")
    slowest_ms = data.get("slowest_tool_ms", 0)
    timeline = data.get("step_timeline", [])

    session_seconds = session_ms / 1000.0

    # Simplified panel when no tool calls (direct LLM response)
    if tool_calls_count == 0:
        return "\n".join([
            "",
            "---",
            f"📊 Session Observability [session={session_id}]",
            "",
            f"🧠 LLM thinking: {total_llm_ms}ms",
            f"⏰ Session time: {session_seconds:.1f}s",
        ])

    # Full panel with tool call details
    lines = [
        "",
        "---",
        f"📊 Session Observability [session={session_id}]",
        "",
        f"🔗 Tool chain: {' → '.join(tool_chain)}",
        f"📈 Tool calls: {tool_calls_count}",
        f"⏱️ Tool time: {total_tool_ms}ms",
        f"🧠 LLM thinking: {total_llm_ms}ms",
        f"⏰ Session time: {session_seconds:.1f}s",
    ]

    if tool_freq:
        freq_parts = [f"{name}×{count}" for name, count in tool_freq.items()]
        lines.append(f"📊 Tool frequency: {', '.join(freq_parts)}")

    if slowest_name:
        lines.append(f"🐢 Slowest tool: {slowest_name} ({slowest_ms}ms)")

    if timeline:
        lines.append("")
        lines.append("📋 Iteration timeline:")
        for idx, step in enumerate(timeline):
            llm_ms = step.get("llm_ms", 0)
            tool_ms = step.get("tool_ms", 0)
            tool_name = step.get("tool_name", "")
            lines.append(f"  Step {idx + 1}: 🧠 {llm_ms}ms → 🔧 {tool_ms}ms ({tool_name})")

    return "\n".join(lines)


async def _non_stream_response(
    agent_runtime: Any,
    session_ctx: dict,
    message: str,
    model: str,
) -> ChatCompletionResponse:
    """生成非流式响应，使用 agent_runtime.run() 和 WebSocket 走相同路径
    
    Args:
        agent_runtime: Agent 运行时
        session_ctx: 会话上下文（包含完整的技能和系统提示）
        message: 用户消息
        model: 模型名称
        
    Returns:
        完整响应
    """
    full_content = ""
    tool_calls_info = []  # 记录工具调用信息
    last_tool_result = None  # 记录最后一个工具的结果
    observability_data = {}  # 性能面板数据
    
    # Use a timeout slightly below the DingTalk client timeout (300s)
    # to ensure we can return a meaningful result before the client
    # gives up.  Import from constants for consistency.
    from pyclaw.constants import AGENT_TIMEOUT_SECONDS as _global_timeout
    AGENT_TIMEOUT_SECONDS = min(_global_timeout, DINGTALK_CLIENT_TIMEOUT)

    async def _consume_agent_events():
        nonlocal full_content, last_tool_result
        async for event in agent_runtime.run(session_ctx, message):
            if event.event_type == "text" and event.content:
                if event.content == "Thinking...":
                    continue
                full_content = event.content
                logger.info("[non-stream] text event: %s...", event.content[:80])
                
            elif event.event_type == "text_delta" and event.content:
                full_content += event.content
                
            elif event.event_type == "tool_call":
                tool_calls_info.append({
                    "tool": event.tool_name,
                    "args": event.tool_args,
                })
                logger.info("[non-stream] Tool call: %s", event.tool_name)
                
            elif event.event_type == "tool_result":
                # 清空 full_content，确保只保留最终一轮（工具执行后）的文本
                full_content = ""
                last_tool_result = event.tool_result
                logger.info("[non-stream] Tool result: %s -> %s...", event.tool_name, str(event.tool_result)[:100])
                
            elif event.event_type == "error":
                logger.error("[non-stream] Agent error: %s", event.content)
                if not full_content:
                    full_content = f"Error: {event.content}"
                    
            elif event.event_type == "done":
                logger.info(
                    "[non-stream] Agent done. full_content_len=%d, tools_used=%d, has_metadata=%s",
                    len(full_content), len(tool_calls_info), bool(event.metadata),
                )
                # 如果 done 事件时 full_content 为空但有工具调用，使用最后一个工具结果作为回复
                if not full_content and tool_calls_info and last_tool_result:
                    # 检查工具结果是否包含有用的信息（如下载链接）
                    result_str = str(last_tool_result)
                    if "http" in result_str.lower() or "下载" in result_str or "link" in result_str.lower():
                        full_content = result_str
                        logger.info("[non-stream] Using last tool result as response: %s...", result_str[:100])
                # Capture observability metadata from done event
                if event.metadata:
                    observability_data.update(event.metadata)

    import time as _time
    session_start_time = _time.monotonic()

    try:
        logger.info("Running agent with full context (skills=%d)", len(session_ctx.get('skills', [])))
        await asyncio.wait_for(_consume_agent_events(), timeout=AGENT_TIMEOUT_SECONDS)
    except asyncio.TimeoutError:
        logger.warning("Agent timed out after %ds, tools_used=%d", AGENT_TIMEOUT_SECONDS, len(tool_calls_info))
        if not full_content:
            # After timeout, if there are tool calls, use the last tool result
            if tool_calls_info and last_tool_result:
                result_str = str(last_tool_result)
                if "http" in result_str.lower() or "下载" in result_str or "link" in result_str.lower():
                    full_content = result_str
                    logger.info("[non-stream] Timeout: using last tool result as response")
                else:
                    tools_desc = ", ".join([t["tool"] for t in tool_calls_info])
                    full_content = f"Task execution timed out ({AGENT_TIMEOUT_SECONDS}s), tools used: {tools_desc}. Please simplify your request and try again."
            else:
                tools_desc = ", ".join([t["tool"] for t in tool_calls_info]) if tool_calls_info else "none"
                full_content = f"Task execution timed out ({AGENT_TIMEOUT_SECONDS}s), tools used: {tools_desc}. Please simplify your request and try again."
        else:
            # Content exists but was truncated by timeout; append a notice
            full_content += "\n\n> ⚠️ *Response may be incomplete due to timeout.*"

        # Build observability data from collected tool_calls_info so the
        # panel is still appended even when the agent did not emit a
        # ``done`` event before the timeout.
        if not observability_data and tool_calls_info:
            elapsed_ms = int((_time.monotonic() - session_start_time) * 1000)
            tool_chain = [t["tool"] for t in tool_calls_info]
            tool_freq: dict[str, int] = {}
            for name in tool_chain:
                tool_freq[name] = tool_freq.get(name, 0) + 1
            observability_data = {
                "session_id": session_ctx.get("session_id", "unknown"),
                "tool_chain": tool_chain,
                "tool_calls_count": len(tool_chain),
                "total_tool_ms": 0,
                "total_llm_ms": elapsed_ms,
                "session_total_ms": elapsed_ms,
                "tool_frequency": tool_freq,
                "slowest_tool": "",
                "slowest_tool_ms": 0,
                "step_timeline": [],
                "timeout": True,
            }
    except Exception as e:
        logger.error("Agent runtime error: %s", e, exc_info=True)
        full_content = f"Sorry, an error occurred while processing your request: {str(e)}"
    
    if not full_content:
        full_content = "Sorry, I am temporarily unable to respond. Please try again later."
    
    # Append Session Observability panel if we have metrics
    if observability_data:
        full_content += _format_observability_panel(observability_data)
    
    # Format Markdown tables for better readability
    full_content = _format_markdown_table(full_content)
    
    return ChatCompletionResponse(
        model=model,
        choices=[
            ChatCompletionChoice(
                message=ChatMessage(role="assistant", content=full_content),
            )
        ],
        usage=ChatCompletionUsage(
            prompt_tokens=len(message) // 4,
            completion_tokens=len(full_content) // 4,
            total_tokens=(len(message) + len(full_content)) // 4,
        ),
    )
